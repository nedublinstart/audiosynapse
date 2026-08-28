from __future__ import annotations

import csv
import io
import re
import zipfile
from html.parser import HTMLParser
from pathlib import Path

from pypdf import PdfReader

IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff", ".heic"}


class _HTMLTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        text = data.strip()
        if text:
            self.parts.append(text)


def extract_text_from_file(path: Path, content_type: str, filename: str) -> str:
    suffix = path.suffix.lower()
    name = filename.lower()

    if suffix == ".pdf" or "pdf" in content_type:
        return _extract_pdf(path)
    if suffix in {".docx"} or "wordprocessingml" in content_type:
        return _extract_docx(path)
    if suffix in {".pptx"} or "presentationml" in content_type or name.endswith(".pptx"):
        return _extract_pptx(path)
    if suffix == ".ppt" or name.endswith(".ppt"):
        return (
            f"[Презентация {filename}] Файл PPT сохранён. "
            "Для полного разбора сохраните как PDF или PPTX."
        )
    if suffix == ".odt" or "opendocument.text" in content_type:
        return _extract_odt(path)
    if suffix == ".rtf":
        return _extract_rtf(path)
    if suffix in {".html", ".htm"} or "html" in content_type:
        return _extract_html(path)
    if suffix == ".csv" or "csv" in content_type:
        return _extract_csv(path)
    if suffix in {".xlsx", ".xls"} or "spreadsheet" in content_type:
        return _extract_xlsx(path)
    if suffix in {".txt", ".md"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    if content_type.startswith("image/") or suffix in IMAGE_SUFFIXES:
        return (
            f"[Изображение: {filename}] Визуальный материал прикреплён. "
            "Если это скан слайда — для текста лучше PDF или DOCX с распознанным текстом."
        )
    return f"[Файл: {filename}] Тип {content_type} сохранён; текстовое извлечение ограничено."


def _extract_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    chunks: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            chunks.append(f"[Слайд: {i}]\n{text}")
    return "\n\n".join(chunks) if chunks else f"[PDF {path.name}] Текст не извлечён — возможно, это скан."


def _extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paras) if paras else f"[DOCX {path.name}] Пустой документ."


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return (
            f"[PPTX {path.name}] Презентация сохранена. "
            "Модуль разбора слайдов недоступен — загрузите PDF с тем же содержимым."
        )

    prs = Presentation(str(path))
    chunks: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                block = (shape.text or "").strip()
                if block:
                    texts.append(block)
        if getattr(slide, "has_notes_slide", False) and slide.notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text.strip():
                texts.append(f"(Заметки докладчика: {notes_frame.text.strip()})")
        if texts:
            chunks.append(f"[Слайд: {i}]\n" + "\n".join(texts))
    return "\n\n".join(chunks) if chunks else f"[PPTX {path.name}] Текст на слайдах не найден."


def _extract_odt(path: Path) -> str:
    try:
        with zipfile.ZipFile(path) as zf:
            xml_bytes = zf.read("content.xml")
    except (OSError, KeyError, zipfile.BadZipFile):
        return f"[ODT {path.name}] Не удалось прочитать документ."

    text = re.sub(r"<text:line-break\s*/>", "\n", xml_bytes.decode("utf-8", errors="ignore"))
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text if text else f"[ODT {path.name}] Текст не извлечён."


def _extract_rtf(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    text = re.sub(r"\\[a-z]+\d*\s?", " ", raw)
    text = re.sub(r"[{}]", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text if len(text) > 40 else f"[RTF {path.name}] Текст не извлечён."


def _extract_html(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    parser = _HTMLTextExtractor()
    parser.feed(raw)
    text = "\n".join(parser.parts)
    return text.strip() if text.strip() else f"[HTML {path.name}] Текст не найден."


def _extract_csv(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    reader = csv.reader(io.StringIO(raw))
    rows = [" | ".join(cell.strip() for cell in row if cell.strip()) for row in reader if any(row)]
    return "\n".join(rows) if rows else f"[CSV {path.name}] Пустой файл."


def _extract_xlsx(path: Path) -> str:
    if path.suffix.lower() == ".xls":
        return (
            f"[XLS {path.name}] Старый формат Excel сохранён. "
            "Сохраните как XLSX или CSV для полного разбора."
        )
    try:
        from openpyxl import load_workbook
    except ImportError:
        return f"[XLSX {path.name}] Таблица сохранена; модуль разбора недоступен."

    wb = load_workbook(str(path), read_only=True, data_only=True)
    chunks: list[str] = []
    for sheet in wb.worksheets:
        lines: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                lines.append(" | ".join(cells))
        if lines:
            chunks.append(f"[Лист: {sheet.title}]\n" + "\n".join(lines[:200]))
    wb.close()
    return "\n\n".join(chunks) if chunks else f"[XLSX {path.name}] Пустая таблица."
